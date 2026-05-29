"""Stage 7: assembly.

- Embeds figures (rendered image or visible placeholder) into the draft.
- Runs a lightweight publication-checklist gate (structural PASS/FAIL).
- Exports .docx via pandoc (Pages opens .docx natively).
- Builds a 2x2 four-panel .pptx (python-pptx) -- replaces the slide + Codia step.
- Optionally opens the .docx in Pages via osascript.
"""
from __future__ import annotations

import shutil
import subprocess
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from . import paths
from .figures import Figure, _FIGURE_RE


# ---------------------------------------------------------------- standard footer
# Heading used as the idempotency marker; appended once before the gate runs.
AUDIENCE_DISCLAIMER_HEADING = "## 対象読者と免責事項"
AUDIENCE_DISCLAIMER_FOOTER = (
    f"\n\n{AUDIENCE_DISCLAIMER_HEADING}\n\n"
    "**対象読者**: 本稿は医師・医療従事者向けである。\n\n"
    "**免責事項**: 本稿は情報提供を目的とし、個別の診療を代替するものではない。"
    "実際の診断・治療にあたっては、最新の添付文書・診療ガイドラインを優先すること。\n"
)


def append_footer(md: str) -> str:
    """Append the audience/disclaimer footer once (idempotent)."""
    if AUDIENCE_DISCLAIMER_HEADING in md:
        return md
    return md.rstrip() + "\n" + AUDIENCE_DISCLAIMER_FOOTER


# ---------------------------------------------------------------- figure embed
def embed_figures(draft_md: str, figures: list[Figure]) -> str:
    blocks = list(_FIGURE_RE.finditer(draft_md))
    if not blocks:
        return draft_md
    by_index = {f.index: f for f in figures}
    out = []
    last = 0
    for i, match in enumerate(blocks, 1):
        out.append(draft_md[last : match.start()])
        fig = by_index.get(i)
        if fig and fig.image_path:
            rel = "../" + fig.image_path  # content/drafts -> repo root
            out.append(f"\n![{fig.title}]({rel})\n\n{fig.caption_ja}\n")
        elif fig:
            caption_q = fig.caption_ja.replace("\n", "\n> ")
            out.append(
                f"\n> **[図{fig.index} ここに挿入]** {fig.title}\n>\n"
                f"> {caption_q}\n>\n"
                f"> _生成プロンプト:_ `{fig.prompt[:120]}...`\n"
            )
        else:
            out.append(match.group(0))
        last = match.end()
    out.append(draft_md[last:])
    return "".join(out)


# ------------------------------------------------------------- publication gate
@dataclass
class GateResult:
    passed: bool
    checks: list[tuple[str, str, str]]  # (item, PASS/FAIL/N/A, comment)

    def report_md(self, title: str) -> str:
        verdict = "✅ 公開可" if self.passed else "❌ 要修正"
        lines = [
            f"# 📋 Publication Check Report: {title}",
            "",
            f"**判定**: {verdict}",
            "",
            "| 項目 | 判定 | コメント |",
            "|------|------|----------|",
        ]
        for item, status, comment in self.checks:
            lines.append(f"| {item} | {status} | {comment} |")
        fails = [c for c in self.checks if c[1] == "FAIL"]
        if fails:
            lines += ["", "## 要修正事項 (FAIL)", ""]
            for item, _, comment in fails:
                lines.append(f"- **{item}**: {comment}")
        return "\n".join(lines) + "\n"


_DISCLAIMER_HINTS = ["免責", "情報提供を目的", "診療の代替", "主治医", "医師・医療従事者向け"]
_COI_HINTS = ["COI", "利益相反", "開示"]


def publication_gate(text: str) -> GateResult:
    """Automated structural subset of content/assets/publication_checklist.md.

    Catches the mechanical FAILs (missing references / disclaimer / citations).
    Substantive medical judgment is still left to the human/LLM reviewer.
    """
    checks: list[tuple[str, str, str]] = []

    has_refs = "## 参考文献" in text
    checks.append(
        ("参考文献セクション", "PASS" if has_refs else "FAIL", "章末に Vancouver 形式の参考文献が必要")
    )

    has_citation = "[" in text and "]" in text and any(f"[{n}]" in text for n in range(1, 60))
    checks.append(
        ("本文中の出典番号 [n]", "PASS" if has_citation else "FAIL", "医学的主張に [n] を付与")
    )

    has_disclaimer = any(h in text for h in _DISCLAIMER_HINTS)
    checks.append(
        ("免責/対象読者の明示", "PASS" if has_disclaimer else "FAIL", "末尾に免責事項 or 対象読者の明記が必要")
    )

    has_coi = any(h in text for h in _COI_HINTS)
    checks.append(
        ("COI開示", "PASS" if has_coi else "N/A", "該当する企業関係がある場合のみ必須")
    )

    not_dry = "(dry/手動" not in text and "provider: none" not in text
    checks.append(
        (
            "本文が生成済み",
            "PASS" if not_dry else "FAIL",
            "LLM 生成済みの本文を検出" if not_dry else "dry モードのプレースホルダのまま。LLM 生成が必要",
        )
    )

    passed = all(s != "FAIL" for _, s, _ in checks)
    return GateResult(passed=passed, checks=checks)


# --------------------------------------------------------------------- exports
def export_docx(md_path: Path, out_dir: Path) -> Path | None:
    if not shutil.which("pandoc"):
        print("[assemble] pandoc not found; skipping .docx export")
        return None
    docx = out_dir / (md_path.stem + ".docx")
    subprocess.run(
        ["pandoc", str(md_path), "-o", str(docx), "--resource-path", str(paths.ROOT)],
        check=True,
    )
    print(f"[assemble] docx -> {docx} (Pages で直接開けます)")
    return docx


def build_pptx(title: str, figures: list[Figure], out_dir: Path) -> Path:
    """4-panel (2x2) deck. Drops the rendered image if present, else empty labeled panels."""
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]

    margin = Emu(457200)
    gap = Emu(152400)
    title_h = Emu(685800)
    grid_top = margin + title_h
    grid_w = prs.slide_width - 2 * margin
    grid_h = prs.slide_height - grid_top - margin
    cell_w = (grid_w - gap) // 2
    cell_h = (grid_h - gap) // 2
    letters = ["A", "B", "C", "D"]

    for fig in figures or [Figure(index=1, title=title)]:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(margin, margin, grid_w, title_h)
        p = tb.text_frame.paragraphs[0]
        p.text = f"図{fig.index}. {fig.title}" if fig.title else title
        p.font.size = Pt(24)
        p.font.bold = True

        for idx, letter in enumerate(letters):
            row, col = divmod(idx, 2)
            x = margin + col * (cell_w + gap)
            y = grid_top + row * (cell_h + gap)
            if fig.image_path:
                img = paths.ROOT / fig.image_path
                if img.exists():
                    # Single rendered 4-panel image fills the grid once.
                    if idx == 0:
                        slide.shapes.add_picture(str(img), margin, grid_top, grid_w, grid_h)
                    continue
            box = slide.shapes.add_shape(1, x, y, cell_w, cell_h)  # rectangle
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            box.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
            lp = box.text_frame.paragraphs[0]
            lp.text = letter
            lp.alignment = PP_ALIGN.LEFT
            lp.font.size = Pt(20)
            lp.font.bold = True
            lp.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        if fig.caption_ja:
            cap = slide.notes_slide.notes_text_frame
            cap.text = f"{fig.caption_ja}\n\nPROMPT: {fig.prompt}"

    pptx_path = out_dir / (paths.slugify(title) + "_figures.pptx")
    prs.save(str(pptx_path))
    print(f"[assemble] pptx -> {pptx_path} (4パネルレイアウト; slide+Codia の置き換え)")
    return pptx_path


def open_in_pages(docx_path: Path) -> None:
    if not shutil.which("osascript") or docx_path is None:
        return
    try:
        subprocess.run(
            ["osascript", "-e", f'tell application "Pages" to open POSIX file "{docx_path}"'],
            check=False,
        )
        print(f"[assemble] Pages で開きました: {docx_path}")
    except Exception as exc:  # noqa: BLE001
        print(f"[assemble] Pages 起動スキップ: {exc}")


def assemble(
    draft_path: Path,
    figures: list[Figure],
    out_dir: Path,
    title: str,
    make_pptx: bool = True,
    open_pages: bool = False,
) -> dict:
    draft_md = draft_path.read_text(encoding="utf-8")
    final_md = embed_figures(draft_md, figures)
    # Append the audience/disclaimer footer before the gate so it lands in the
    # final markdown AND satisfies the 免責/対象読者 check (idempotent on re-runs).
    final_md = append_footer(final_md)

    final_path = out_dir / (draft_path.stem + "_final.md")
    final_path.write_text(final_md, encoding="utf-8")

    gate = publication_gate(final_md)
    gate_path = out_dir / "publication_check.md"
    gate_path.write_text(gate.report_md(title), encoding="utf-8")
    print(f"[assemble] publication gate: {'PASS' if gate.passed else 'FAIL'} -> {gate_path}")

    docx = export_docx(final_path, out_dir)
    pptx = build_pptx(title, figures, out_dir) if make_pptx else None

    if open_pages and docx:
        open_in_pages(docx)

    return {
        "final_md": str(final_path),
        "docx": str(docx) if docx else None,
        "pptx": str(pptx) if pptx else None,
        "gate_passed": gate.passed,
        "gate_report": str(gate_path),
    }
